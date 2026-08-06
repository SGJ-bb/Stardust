"""星尘 Stradust 项目演示答辩 PPT 生成器"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色定义 ──
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)       # 深色背景
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)       # 卡片背景
ACCENT_PURPLE = RGBColor(0x8B, 0x6C, 0xFF)  # 主强调紫
ACCENT_BLUE = RGBColor(0x5E, 0xA2, 0xFF)    # 辅助蓝
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xB0)   # 绿色
ACCENT_YELLOW = RGBColor(0xFF, 0xC8, 0x5E)  # 黄色
ACCENT_PINK = RGBColor(0xFF, 0x7E, 0xB3)    # 粉色
TEXT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)     # 白色文字
TEXT_MUTED = RGBColor(0xA0, 0xA0, 0xB0)     # 辅助文字
TEXT_DIM = RGBColor(0x6C, 0x6C, 0x80)       # 暗文字
GRADIENT_START = RGBColor(0x8B, 0x6C, 0xFF)
GRADIENT_END = RGBColor(0x5E, 0xA2, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BG_DARK):
    """添加纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=TEXT_WHITE, bullet_color=ACCENT_PURPLE):
    """添加带项目符号的列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(8)
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_PURPLE):
    """添加装饰线"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_icon_card(slide, left, top, width, height, icon, title, desc, accent=ACCENT_PURPLE):
    """添加图标卡片"""
    card = add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40), border_width=Pt(1))
    # 左侧装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.15), Pt(4), height - Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # 图标
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.6), Inches(0.5), icon, font_size=24, color=accent, alignment=PP_ALIGN.CENTER)
    # 标题
    add_text_box(slide, left + Inches(0.85), top + Inches(0.1), width - Inches(1.1), Inches(0.35), title, font_size=16, color=TEXT_WHITE, bold=True)
    # 描述
    add_text_box(slide, left + Inches(0.85), top + Inches(0.45), width - Inches(1.1), height - Inches(0.55), desc, font_size=12, color=TEXT_MUTED)
    return card

# ════════════════════════════════════════════════════════════════
# 第1页：封面
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_bg(slide)

# 装饰圆
circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(5), Inches(5))
circle1.fill.solid()
circle1.fill.fore_color.rgb = RGBColor(0x1A, 0x10, 0x30)
circle1.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(6), Inches(6))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x10, 0x10, 0x25)
circle2.line.fill.background()

# 项目名
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2), "星尘 Stradust", font_size=56, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5.5), Inches(3.0), Inches(2.3), ACCENT_PURPLE)

# 副标题
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.8), "AI 伴侣应用 — 跨平台重构与深度优化", font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# 描述
add_text_box(slide, Inches(2.5), Inches(4.5), Inches(8), Inches(0.6), "基于 Tauri 2 + React 19 + Rust 的 iOS/PC 跨平台实现", font_size=18, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# 底部信息
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.4), "2026年6月  |  毕业设计答辩", font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# 第2页：项目概述
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6), "项目概述", font_size=32, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_PURPLE)

# 左侧：项目定位
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5), "项目定位", font_size=20, color=ACCENT_PURPLE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(1.5),
    "星尘是一款 AI 伴侣应用，核心角色是一只异色瞳黑猫「星尘」。\n"
    "用户可以与 AI 角色进行深度对话、共同写日记、分享朋友圈、\n"
    "管理记忆、培养好感度，体验完整的情感陪伴旅程。",
    font_size=15, color=TEXT_MUTED)

# 右侧：核心能力卡片
add_text_box(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.5), "核心能力", font_size=20, color=ACCENT_BLUE, bold=True)

features = [
    ("🧠", "智能记忆", "5层记忆架构：短期→长期→跨场景→会话→难忘时刻", ACCENT_PURPLE),
    ("💬", "多角色对话", "支持多角色私聊、群聊、微信iLink消息互通", ACCENT_BLUE),
    ("❤️", "好感度系统", "0-100好感度，6级称号，行为评估，性格进化", ACCENT_PINK),
    ("🎮", "游戏化系统", "成就/签到/里程碑/时光胶囊/纪念相册", ACCENT_GREEN),
    ("📝", "日记与朋友圈", "AI生成日记、朋友圈动态、评论互动", ACCENT_YELLOW),
    ("🎨", "个性化定制", "角色创建、皮肤商店、贴纸、气泡样式", ACCENT_PURPLE),
]

for i, (icon, title, desc, accent) in enumerate(features):
    row = i // 2
    col = i % 2
    add_icon_card(slide,
        Inches(7 + col * 3.0), Inches(1.9 + row * 1.35),
        Inches(2.8), Inches(1.2),
        icon, title, desc, accent)

# 底部：技术栈
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40), border_width=Pt(1))
add_text_box(slide, Inches(1.0), Inches(5.9), Inches(2), Inches(0.4), "技术栈", font_size=16, color=ACCENT_PURPLE, bold=True)

techs = [
    ("Android 原版", "Kotlin + OkHttp + TFLite + Live2D", ACCENT_GREEN),
    ("iOS 版", "Tauri 2 + React 19 + Rust + Anime.js", ACCENT_BLUE),
    ("PC 版", "Tauri 2 + React 19 + Rust + Electron", ACCENT_PURPLE),
    ("后端", "Python FastAPI + OpenAI + Mem0", ACCENT_YELLOW),
]
for i, (name, tech, color) in enumerate(techs):
    x = Inches(1.0 + i * 3.0)
    add_text_box(slide, x, Inches(6.25), Inches(2.8), Inches(0.3), name, font_size=13, color=color, bold=True)
    add_text_box(slide, x, Inches(6.55), Inches(2.8), Inches(0.3), tech, font_size=11, color=TEXT_DIM)

# ════════════════════════════════════════════════════════════════
# 第3页：系统架构
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6), "系统架构", font_size=32, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_PURPLE)

# Android 架构
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5), "Android 原版架构", font_size=20, color=ACCENT_GREEN, bold=True)

layers_android = [
    ("UI 层", "32个Activity + ChatViewModel + Adapter", ACCENT_GREEN),
    ("服务层", "5个前台Service + 3个BroadcastReceiver", ACCENT_BLUE),
    ("核心引擎层", "ContextManager + MemoryPool + SessionManager", ACCENT_PURPLE),
    ("中间件层", "ApiClient + PromptBuilder + PluginRegistry + RAG", ACCENT_YELLOW),
    ("存储层", "SharedPreferences(加密) + ChatHistoryStorage + 文件系统", TEXT_DIM),
]
for i, (name, desc, color) in enumerate(layers_android):
    y = Inches(1.85 + i * 0.65)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Pt(4), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text_box(slide, Inches(1.0), y, Inches(1.5), Inches(0.5), name, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(2.5), y, Inches(4), Inches(0.5), desc, font_size=13, color=TEXT_MUTED)

# iOS 架构
add_text_box(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5), "iOS 版架构 (Tauri 2)", font_size=20, color=ACCENT_BLUE, bold=True)

layers_ios = [
    ("前端层", "React 19 + TypeScript + Vite + Anime.js", ACCENT_BLUE),
    ("16个页面", "Chat/Diary/Achievement/Memory/Moments/GroupChat/...", ACCENT_GREEN),
    ("后端层", "Rust + reqwest + serde + chrono + regex", ACCENT_PURPLE),
    ("29个命令", "send_chat / generate_diary / evaluate_memories / ...", ACCENT_YELLOW),
    ("存储层", "JSON文件持久化 + app_data_dir 沙盒存储", TEXT_DIM),
]
for i, (name, desc, color) in enumerate(layers_ios):
    y = Inches(1.85 + i * 0.65)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), y, Pt(4), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text_box(slide, Inches(7.2), y, Inches(1.5), Inches(0.5), name, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(8.7), y, Inches(4), Inches(0.5), desc, font_size=13, color=TEXT_MUTED)

# 数据流
add_shape(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.8), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40), border_width=Pt(1))
add_text_box(slide, Inches(1.0), Inches(5.4), Inches(4), Inches(0.4), "核心数据流：私聊消息", font_size=16, color=ACCENT_PURPLE, bold=True)

flow_items = [
    "用户输入", "→", "Prompt构建", "→", "记忆注入", "→", "RAG检索", "→", "LLM调用", "→", "工具执行", "→", "记忆提取", "→", "好感度更新"
]
x = Inches(1.0)
for item in flow_items:
    if item == "→":
        add_text_box(slide, x, Inches(5.9), Inches(0.4), Inches(0.4), "→", font_size=16, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
        x += Inches(0.4)
    else:
        w = Inches(1.2)
        box = add_shape(slide, x, Inches(5.85), w, Inches(0.4), fill_color=RGBColor(0x25, 0x25, 0x40), border_color=ACCENT_PURPLE, border_width=Pt(1))
        add_text_box(slide, x, Inches(5.85), w, Inches(0.4), item, font_size=11, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
        x += Inches(1.25)

add_text_box(slide, Inches(1.0), Inches(6.4), Inches(11), Inches(0.4),
    "Android: ChatViewModel → PromptBuilder → ContextManager → PersonaRagManager → ApiClient → PluginRegistry",
    font_size=12, color=TEXT_DIM)
add_text_box(slide, Inches(1.0), Inches(6.7), Inches(11), Inches(0.4),
    "iOS: React State → invoke('send_chat') → Rust reqwest → OpenAI API → extract_emotion → invoke callback",
    font_size=12, color=TEXT_DIM)

# ════════════════════════════════════════════════════════════════
# 第4页：记忆系统
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6), "记忆系统 — 核心架构", font_size=32, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_PURPLE)

# 5层记忆架构
memory_layers = [
    ("MemorableMomentsManager", "难忘时刻", "高分记忆 (score>=8)", Inches(0.8), ACCENT_PINK),
    ("SessionManager", "会话管理", "会话轮换，MAX 20个会话", Inches(1.7), ACCENT_YELLOW),
    ("MemoryPool", "场景记忆池", "总结记忆+细节记忆，MAX 3000字", Inches(2.6), ACCENT_GREEN),
    ("GlobalMemoryPool", "跨场景共享记忆", "isGlobal=true的记忆，跨私聊/群聊", Inches(3.5), ACCENT_BLUE),
    ("MemoryManager", "长期记忆", "简单事实存储，最多200条", Inches(4.4), ACCENT_PURPLE),
]

for name, cn, desc, y, color in memory_layers:
    add_shape(slide, Inches(1.0), y, Inches(5.5), Inches(0.75), fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(2.5), Inches(0.35), name, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.38), Inches(2.5), Inches(0.3), cn, font_size=12, color=TEXT_WHITE)
    add_text_box(slide, Inches(3.8), y + Inches(0.15), Inches(2.5), Inches(0.45), desc, font_size=12, color=TEXT_MUTED)

# 协调器
add_shape(slide, Inches(1.0), Inches(5.35), Inches(5.5), Inches(0.6), fill_color=RGBColor(0x20, 0x18, 0x35), border_color=ACCENT_PURPLE, border_width=Pt(2))
add_text_box(slide, Inches(1.2), Inches(5.4), Inches(5), Inches(0.5), "ContextManager — 核心协调器", font_size=16, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# 右侧：记忆提取流程
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.5), "记忆提取流程", font_size=20, color=ACCENT_BLUE, bold=True)

steps = [
    "1. 对话轮次累积 → shouldEvaluate() (每N轮)",
    "2. evaluateAndUpdateMemory() → AI评估对话",
    "3. MemoryPool.evaluateTurn() → 提取结构化记忆JSON",
    "4. 解析结果 → add/update/delete/detail 操作",
    "5. GlobalMemoryPool.addFromScene() (跨场景)",
    "6. MemoryPool.consolidate() (每10轮整合压缩)",
    "7. SessionManager.checkMemoryLimit() → 超限创建新会话",
]
for i, step in enumerate(steps):
    y = Inches(1.9 + i * 0.5)
    color = ACCENT_BLUE if i % 2 == 0 else TEXT_MUTED
    add_text_box(slide, Inches(7.4), y, Inches(5.3), Inches(0.45), step, font_size=13, color=color)

# iOS 实现
add_shape(slide, Inches(7.2), Inches(5.5), Inches(5.5), Inches(1.5), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40), border_width=Pt(1))
add_text_box(slide, Inches(7.4), Inches(5.6), Inches(5), Inches(0.4), "iOS 版实现", font_size=16, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(7.4), Inches(6.0), Inches(5), Inches(0.8),
    "• Rust 后端 evaluate_memories 命令\n"
    "• AI 评估对话 → 提取记忆 → 分类(habit/preference/impression/detail)\n"
    "• 前端 MemoryPage 可视化管理 + 分类筛选 + 长按删除",
    font_size=12, color=TEXT_MUTED)

# ════════════════════════════════════════════════════════════════
# 第5页：好感度与游戏化
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6), "好感度与游戏化系统", font_size=32, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_PINK)

# 好感度等级
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5), "好感度等级体系", font_size=20, color=ACCENT_PINK, bold=True)

levels = [
    ("0-20", "陌生人", RGBColor(0xFF, 0x45, 0x45)),
    ("21-40", "认识的人", RGBColor(0xFF, 0x7E, 0x45)),
    ("41-60", "朋友", RGBColor(0xFF, 0xC8, 0x5E)),
    ("61-80", "好朋友", RGBColor(0x4E, 0xC9, 0xB0)),
    ("81-95", "亲密伙伴", RGBColor(0xFF, 0x7E, 0xB3)),
    ("96-100", "最重要的人", RGBColor(0xFF, 0x7E, 0xB3)),
]
for i, (range_str, title, color) in enumerate(levels):
    y = Inches(1.85 + i * 0.55)
    # 进度条背景
    add_shape(slide, Inches(0.8), y, Inches(5.5), Inches(0.4), fill_color=RGBColor(0x15, 0x15, 0x25))
    # 进度条填充
    pct = int(range_str.split("-")[1]) / 100
    add_shape(slide, Inches(0.8), y, Inches(5.5 * pct), Inches(0.4), fill_color=color)
    add_text_box(slide, Inches(0.9), y, Inches(1), Inches(0.4), range_str, font_size=12, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(2.2), y, Inches(2), Inches(0.4), title, font_size=13, color=TEXT_WHITE, bold=True)

# 行为评估
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(0.5), "行为评估机制", font_size=16, color=ACCENT_PINK, bold=True)
behaviors = [
    "冒犯行为 → -2好感度",
    "HAPPY情绪 → +1好感度",
    "ANGRY情绪 → -1好感度",
    "傲娇反应 → +1好感度",
    "正面词汇 → +1好感度",
    "负面词汇 → -2好感度",
    "好感度≥90 → 不再减少",
    "每增加5点 → 触发性格进化",
]
for i, b in enumerate(behaviors):
    col = i // 4
    row = i % 4
    add_text_box(slide, Inches(0.8 + col * 2.8), Inches(5.75 + row * 0.35), Inches(2.8), Inches(0.3), "• " + b, font_size=11, color=TEXT_MUTED)

# 右侧：游戏化系统
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.5), "游戏化系统", font_size=20, color=ACCENT_GREEN, bold=True)

gamify = [
    ("🏆", "成就系统", "多类别成就追踪，进度计算，自动更新", ACCENT_YELLOW),
    ("📅", "每日签到", "+2好感度/天，连续签到奖励", ACCENT_GREEN),
    ("⭐", "里程碑", "首次对话/100条消息/7天连续等", ACCENT_PURPLE),
    ("⏳", "时光胶囊", "封存消息，定时开启，3种状态", ACCENT_BLUE),
    ("🖼️", "纪念相册", "回忆记录，分类筛选(milestone/first/special)", ACCENT_PINK),
    ("👗", "皮肤商店", "5种默认气泡皮肤，实时预览", ACCENT_YELLOW),
    ("🎨", "贴纸系统", "9种默认贴纸，情绪分类筛选", ACCENT_GREEN),
    ("⏰", "定时唤醒", "自定义唤醒任务，启用/禁用", ACCENT_BLUE),
]
for i, (icon, title, desc, accent) in enumerate(gamify):
    row = i // 2
    col = i % 2
    add_icon_card(slide,
        Inches(7.2 + col * 2.9), Inches(1.85 + row * 1.3),
        Inches(2.7), Inches(1.15),
        icon, title, desc, accent)

# ════════════════════════════════════════════════════════════════
# 第6页：iOS版技术实现
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), "iOS 版技术实现", font_size=32, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_BLUE)

# 技术栈详情
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5), "技术栈", font_size=20, color=ACCENT_BLUE, bold=True)

stack_items = [
    ("框架", "Tauri 2 (Rust 后端 + Web 前端)", ACCENT_PURPLE),
    ("前端", "React 19 + TypeScript + Vite", ACCENT_BLUE),
    ("后端", "Rust (reqwest + serde + chrono + regex)", ACCENT_GREEN),
    ("动画", "Anime.js v4 + 自定义动画工具模块", ACCENT_YELLOW),
    ("设计", "Impeccable 设计原则 + Taste Skill 反模板", ACCENT_PINK),
    ("构建", "GitHub Actions 云构建 + xcodebuild", ACCENT_BLUE),
]
for i, (label, value, color) in enumerate(stack_items):
    y = Inches(1.85 + i * 0.55)
    add_shape(slide, Inches(0.8), y, Inches(1.2), Inches(0.4), fill_color=color)
    add_text_box(slide, Inches(0.8), y, Inches(1.2), Inches(0.4), label, font_size=13, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.1), y, Inches(4), Inches(0.4), value, font_size=14, color=TEXT_MUTED)

# 右侧：16个页面
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.5), "16个功能页面", font_size=20, color=ACCENT_GREEN, bold=True)

pages = [
    ("ChatPage", "核心聊天，情绪标签，打字指示器"),
    ("DiaryPage", "AI生成日记，心情选择器"),
    ("AchievementPage", "成就统计，进度计算，分类展示"),
    ("MemoryPage", "记忆管理，AI评估，分类筛选"),
    ("ProfilePage", "好感度，签到，角色切换"),
    ("MomentsPage", "朋友圈，AI动态，评论互动"),
    ("GroupChatPage", "群聊创建，多角色对话"),
    ("TimeCapsulePage", "时光胶囊，3种状态管理"),
    ("SettingsPage", "API配置，用户设置，测试连接"),
    ("MorePage", "功能入口网格，个人资料卡片"),
    ("SkinShopPage", "皮肤预览，激活切换"),
    ("StickerPage", "贴纸管理，情绪分类"),
    ("WakeUpPage", "定时唤醒任务管理"),
    ("ChatHistoryPage", "聊天记录搜索/筛选/导出"),
    ("CalendarPage", "月历视图，日期指示器"),
    ("MemorialAlbumPage", "回忆记录，分类筛选"),
]
for i, (name, desc) in enumerate(pages):
    row = i // 2
    col = i % 2
    x = Inches(7.2 + col * 2.9)
    y = Inches(1.85 + row * 0.6)
    add_text_box(slide, x, y, Inches(2.8), Inches(0.3), name, font_size=12, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, x, y + Inches(0.25), Inches(2.8), Inches(0.3), desc, font_size=10, color=TEXT_DIM)

# Rust 后端
add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.6), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40), border_width=Pt(1))
add_text_box(slide, Inches(1.0), Inches(5.6), Inches(4), Inches(0.4), "Rust 后端 — 29个 Tauri 命令", font_size=16, color=ACCENT_PURPLE, bold=True)

rust_cmds = [
    "send_chat / test_connection / save_settings / load_settings",
    "save_chat_history / load_chat_history / save_affection / load_affection",
    "save_achievements / load_achievements / save_checkin_records / load_checkin_records",
    "save_diaries / load_diaries / save_memories / load_memories",
    "save_personas / load_personas / save_time_capsules / load_time_capsules",
    "save_moments / load_moments / save_group_chats / load_group_chats",
    "generate_diary / generate_proactive_chat / generate_moment",
    "evaluate_memories / evolve_personality / web_search",
]
for i, cmd in enumerate(rust_cmds):
    add_text_box(slide, Inches(1.0 + (i % 2) * 5.8), Inches(6.0 + (i // 2) * 0.3), Inches(5.6), Inches(0.3), cmd, font_size=10, color=TEXT_DIM)

# ════════════════════════════════════════════════════════════════
# 第7页：UI设计系统
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), "UI 设计系统 v2", font_size=32, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_PURPLE)

# 设计系统6大模块
design_modules = [
    ("排版系统", [
        "模块化字号比例 (1.25 Major Third)",
        "4级字重: regular/medium/semibold/bold",
        "3级行高/字间距: tight/normal/wide",
    ], ACCENT_PURPLE),
    ("色彩系统", [
        "紫调中性色 (禁止纯灰)",
        "7色强调色板 + 4组渐变",
        "4级阴影: sm/md/lg/glow",
    ], ACCENT_BLUE),
    ("空间系统", [
        "4px基准网格 (10级空间)",
        "6级圆角: xs~2xl",
        "大/中/小节奏交替",
    ], ACCENT_GREEN),
    ("动效系统", [
        "4条缓动曲线 (禁止bounce)",
        "5级时长: instant~sluggish",
        "3级交错延迟",
    ], ACCENT_YELLOW),
    ("交互系统", [
        "44px最小触摸目标",
        "focus发光环 + active缩放",
        "prefers-reduced-motion 支持",
    ], ACCENT_PINK),
    ("动画模块", [
        "20+ Anime.js 动画函数",
        "fadeInUp/fadeInScale/staggerFadeIn",
        "messageBubbleIn/bottomSheet/breathe",
    ], ACCENT_PURPLE),
]

for i, (title, items, color) in enumerate(design_modules):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.3 + row * 2.8)
    add_shape(slide, x, y, Inches(3.8), Inches(2.5), fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.4), title, font_size=18, color=color, bold=True)
    add_accent_line(slide, x + Inches(0.2), y + Inches(0.55), Inches(1.5), color)
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.2), y + Inches(0.7 + j * 0.4), Inches(3.4), Inches(0.35), "• " + item, font_size=12, color=TEXT_MUTED)

# ════════════════════════════════════════════════════════════════
# 第8页：功能覆盖度对比
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), "功能覆盖度对比", font_size=32, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_GREEN)

# 已实现
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5), "已实现功能 (77%)", font_size=20, color=ACCENT_GREEN, bold=True)

implemented = [
    "核心聊天 + 情绪标签 + 打字指示器",
    "多角色系统 (创建/切换/编辑/删除)",
    "5层记忆架构 + AI评估记忆",
    "好感度系统 (0-100, 6级称号, 行为评估)",
    "日记系统 (AI生成, 心情选择)",
    "成就系统 + 每日签到 + 里程碑",
    "朋友圈 + AI动态 + 评论互动",
    "群聊系统 (多角色对话, 链式触发)",
    "时光胶囊 + 纪念相册 + 日历",
    "皮肤商店 + 贴纸系统",
    "定时唤醒 + 聊天记录管理",
    "角色性格进化 + AI主动搭话",
    "消息反应 + 人性化处理",
    "内容安全过滤 + 情感守护",
    "聊天预测 + 网页搜索",
    "设置管理 + API配置 + 测试连接",
]
for i, item in enumerate(implemented):
    row = i // 2
    col = i % 2
    add_text_box(slide, Inches(0.8 + col * 3.0), Inches(1.8 + row * 0.4), Inches(2.9), Inches(0.35), "✓ " + item, font_size=11, color=ACCENT_GREEN)

# 未实现
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.5), "未实现功能 (23%)", font_size=20, color=ACCENT_YELLOW, bold=True)

not_implemented = [
    ("Live2D 模型系统", "需 iOS 原生 WebView 插件", ACCENT_YELLOW),
    ("语音通话 / TTS / ASR", "需 iOS 原生音频插件", ACCENT_YELLOW),
    ("睡前电台", "依赖 TTS 引擎", ACCENT_YELLOW),
    ("微信 iLink 连接", "iOS 无后台轮询能力", RGBColor(0xFF, 0x60, 0x60)),
    ("悬浮窗/屏幕识别", "iOS 无此权限", RGBColor(0xFF, 0x60, 0x60)),
    ("本地模型/TFLite", "需 CoreML 适配", ACCENT_YELLOW),
    ("虚拟世界", "需独立渲染引擎", ACCENT_YELLOW),
    ("RAG 向量检索", "需 iOS 端嵌入模型", ACCENT_YELLOW),
]
for i, (name, reason, color) in enumerate(not_implemented):
    y = Inches(1.85 + i * 0.6)
    add_shape(slide, Inches(7.2), y, Inches(5.3), Inches(0.5), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40), border_width=Pt(1))
    add_text_box(slide, Inches(7.4), y + Inches(0.05), Inches(2.5), Inches(0.4), "✗ " + name, font_size=12, color=color, bold=True)
    add_text_box(slide, Inches(9.9), y + Inches(0.05), Inches(2.5), Inches(0.4), reason, font_size=11, color=TEXT_DIM)

# 底部说明
add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40), border_width=Pt(1))
add_text_box(slide, Inches(1.0), Inches(6.55), Inches(11), Inches(0.6),
    "注：红色标记为 iOS 平台限制导致无法实现的功能，黄色标记为需要原生插件支持的功能，可在后续版本中逐步补全。",
    font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# 第9页：Bug修复与质量保障
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), "Bug 修复与质量保障", font_size=32, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_YELLOW)

# 关键 Bug 修复
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5), "关键 Bug 修复", font_size=20, color=ACCENT_YELLOW, bold=True)

bugs = [
    ("Rust 参数映射", "_persona_name 前缀下划线破坏 Tauri 2 映射", "改为 persona_name + #[allow(unused_variables)]", ACCENT_RED if hasattr(ACCENT_RED := RGBColor(0xFF, 0x45, 0x45), 'rgb') else RGBColor(0xFF, 0x45, 0x45)),
    ("角色切换无效", "ChatPage 硬编码 DEFAULT_CHARACTER.system_prompt", "加载活跃角色的 system_prompt", RGBColor(0xFF, 0x45, 0x45)),
    ("分类筛选不匹配", "前端中文/Rust英文导致 AI 评估记忆无法筛选", "统一使用英文 key + 中文映射", RGBColor(0xFF, 0x45, 0x45)),
    ("async/await 错误", "MemorialAlbumPage 同步调用 async 函数", "改为 IIFE async 模式", RGBColor(0xFF, 0x80, 0x00)),
    ("sort 原地修改", "4个页面 sort() 原地修改传入数组", "改为 [...arr].sort() 先复制", RGBColor(0xFF, 0x80, 0x00)),
    ("动画内存泄漏", "7处动画缺少 cleanup 导致内存泄漏", "添加 useEffect cleanup + revert()", RGBColor(0xFF, 0x80, 0x00)),
    ("CSS 变量不存在", "var(--accent-cool) 引用不存在的变量", "改为 var(--accent-secondary)", ACCENT_YELLOW),
    ("圆角变量错误", "borderBottomLeftRadius 用了空间变量", "改为 var(--radius-sm)", ACCENT_YELLOW),
]
for i, (name, problem, fix, color) in enumerate(bugs):
    y = Inches(1.8 + i * 0.6)
    add_shape(slide, Inches(0.8), y, Inches(5.5), Inches(0.5), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40), border_width=Pt(1))
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(1.5), Inches(0.4), name, font_size=11, color=color, bold=True)
    add_text_box(slide, Inches(2.5), y + Inches(0.05), Inches(1.8), Inches(0.4), problem, font_size=10, color=TEXT_MUTED)
    add_text_box(slide, Inches(4.3), y + Inches(0.05), Inches(1.8), Inches(0.4), "→ " + fix, font_size=10, color=ACCENT_GREEN)

# 右侧：质量保障
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.5), "质量保障措施", font_size=20, color=ACCENT_GREEN, bold=True)

qa_items = [
    ("编译验证", "Rust cargo check: 0 errors\nVite build: 98 modules 通过", ACCENT_GREEN),
    ("代码审查", "全部16个页面文件逐行审查\napi.ts + lib.rs + App.tsx 深度检查", ACCENT_BLUE),
    ("Bug 修复", "2轮全面复查，共修复 15+ bug\n包括参数映射/内存泄漏/逻辑错误", ACCENT_YELLOW),
    ("功能对比", "Android 28类功能逐一对比\n确认77%覆盖度，23%为平台限制", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(qa_items):
    y = Inches(1.85 + i * 1.2)
    add_shape(slide, Inches(7.2), y, Inches(5.3), Inches(1.0), fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    add_text_box(slide, Inches(7.4), y + Inches(0.1), Inches(2), Inches(0.35), title, font_size=15, color=color, bold=True)
    add_text_box(slide, Inches(7.4), y + Inches(0.45), Inches(5), Inches(0.5), desc, font_size=11, color=TEXT_MUTED)

# ════════════════════════════════════════════════════════════════
# 第10页：CI/CD 与部署
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), "CI/CD 与部署方案", font_size=32, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_BLUE)

# GitHub Actions 流程
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5), "GitHub Actions 云构建流程", font_size=20, color=ACCENT_BLUE, bold=True)

ci_steps = [
    ("1", "Checkout", "拉取代码", ACCENT_PURPLE),
    ("2", "Setup", "Node.js 22 + pnpm 9 + Rust stable", ACCENT_BLUE),
    ("3", "Install", "前端依赖 + Rust iOS targets", ACCENT_GREEN),
    ("4", "Init", "tauri ios init (占位Team ID)", ACCENT_YELLOW),
    ("5", "Build", "tauri ios build --debug (跳过签名)", ACCENT_PINK),
    ("6", "Package", "查找 .app → 打包为 .ipa", ACCENT_PURPLE),
    ("7", "Upload", "上传 Artifact (保留30天)", ACCENT_BLUE),
]
for i, (num, name, desc, color) in enumerate(ci_steps):
    y = Inches(1.85 + i * 0.6)
    # 步骤编号圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05), Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(0.35), Inches(0.35), num, font_size=12, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.3), y + Inches(0.05), Inches(1.5), Inches(0.35), name, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(2.8), y + Inches(0.05), Inches(3.5), Inches(0.35), desc, font_size=12, color=TEXT_MUTED)

# 右侧：部署方案
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.5), "部署方案", font_size=20, color=ACCENT_GREEN, bold=True)

deploy = [
    ("Debug 构建", "无需 Apple Developer 账号\n通过 AltStore/Sideloadly 安装\n7天有效期（免费Apple ID签名）", ACCENT_GREEN),
    ("Release 构建", "需 Apple Developer 账号 ($99/年)\n配置证书+描述文件到 GitHub Secrets\nAd Hoc / App Store 分发", ACCENT_BLUE),
    ("TestFlight", "通过 App Store Connect 上传\n邀请测试用户（最多10000人）\n正式发布前的最佳测试方案", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(deploy):
    y = Inches(1.85 + i * 1.6)
    add_shape(slide, Inches(7.2), y, Inches(5.3), Inches(1.4), fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    add_text_box(slide, Inches(7.4), y + Inches(0.1), Inches(5), Inches(0.35), title, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(7.4), y + Inches(0.45), Inches(5), Inches(0.9), desc, font_size=12, color=TEXT_MUTED)

# ════════════════════════════════════════════════════════════════
# 第11页：项目亮点与创新
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), "项目亮点与创新", font_size=32, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_YELLOW)

highlights = [
    ("跨平台架构", "使用 Tauri 2 实现一套代码同时支持 iOS/PC/Android，Rust 后端保证性能和安全，Web 前端保证开发效率和UI一致性", ACCENT_PURPLE),
    ("5层记忆架构", "从短期对话到长期事实的5层记忆体系，AI自动评估提取记忆，跨场景共享，会话轮换，记忆整合压缩", ACCENT_BLUE),
    ("深度情感系统", "0-100好感度 + 6级称号 + 行为评估 + 性格进化，AI根据用户行为动态调整角色性格和互动方式", ACCENT_PINK),
    ("AI原生设计", "AI生成日记/朋友圈/主动搭话，AI评估记忆重要性，AI进化角色性格，AI驱动的群聊多角色对话", ACCENT_GREEN),
    ("高品质UI", "Impeccable设计原则 + Anime.js动画 + Taste Skill反模板，20+动画函数，6大设计系统模块", ACCENT_YELLOW),
    ("完整游戏化", "成就/签到/里程碑/时光胶囊/纪念相册/皮肤/贴纸，7个游戏化子系统深度整合", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(highlights):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.3 + row * 1.9)
    add_shape(slide, x, y, Inches(5.8), Inches(1.7), fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    # 左侧装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.15), Pt(5), Inches(1.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(5.3), Inches(0.4), title, font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.6), Inches(5.3), Inches(1.0), desc, font_size=13, color=TEXT_MUTED)

# ════════════════════════════════════════════════════════════════
# 第12页：未来展望
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6), "未来展望", font_size=32, color=TEXT_WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2), ACCENT_GREEN)

# 短期
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(3.7), Inches(0.5), "短期目标", font_size=20, color=ACCENT_GREEN, bold=True)
short_term = [
    "CoreML 适配本地模型推理",
    "AVFoundation 语音通话/TTS/ASR",
    "推送通知 + 后台任务",
    "App Store 上架准备",
]
for i, item in enumerate(short_term):
    add_text_box(slide, Inches(0.8), Inches(1.85 + i * 0.45), Inches(3.7), Inches(0.4), "▸ " + item, font_size=14, color=TEXT_MUTED)

# 中期
add_text_box(slide, Inches(4.8), Inches(1.3), Inches(3.7), Inches(0.5), "中期目标", font_size=20, color=ACCENT_BLUE, bold=True)
mid_term = [
    "Live2D iOS 原生渲染",
    "RAG 向量检索 (CoreML嵌入)",
    "多语言国际化 (i18n)",
    "iCloud 数据同步",
]
for i, item in enumerate(mid_term):
    add_text_box(slide, Inches(4.8), Inches(1.85 + i * 0.45), Inches(3.7), Inches(0.4), "▸ " + item, font_size=14, color=TEXT_MUTED)

# 长期
add_text_box(slide, Inches(8.8), Inches(1.3), Inches(3.7), Inches(0.5), "长期愿景", font_size=20, color=ACCENT_PURPLE, bold=True)
long_term = [
    "AR 虚拟世界 (ARKit)",
    "Apple Watch 伴侣应用",
    "Siri 快捷指令集成",
    "跨设备情感同步",
]
for i, item in enumerate(long_term):
    add_text_box(slide, Inches(8.8), Inches(1.85 + i * 0.45), Inches(3.7), Inches(0.4), "▸ " + item, font_size=14, color=TEXT_MUTED)

# 技术路线图
add_shape(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(3.0), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2A, 0x40), border_width=Pt(1))
add_text_box(slide, Inches(1.0), Inches(4.1), Inches(4), Inches(0.4), "技术路线图", font_size=18, color=ACCENT_YELLOW, bold=True)

roadmap = [
    ("v1.0", "当前版本", "16个页面 + 29个Rust命令 + 设计系统v2", ACCENT_GREEN),
    ("v1.1", "语音版", "AVFoundation语音 + 推送通知 + 后台任务", ACCENT_BLUE),
    ("v1.2", "智能版", "CoreML本地模型 + RAG检索 + 性格进化增强", ACCENT_PURPLE),
    ("v2.0", "沉浸版", "Live2D + AR虚拟世界 + Apple Watch + Siri", ACCENT_YELLOW),
]
for i, (ver, name, desc, color) in enumerate(roadmap):
    y = Inches(4.6 + i * 0.55)
    add_shape(slide, Inches(1.0), y, Inches(0.8), Inches(0.4), fill_color=color)
    add_text_box(slide, Inches(1.0), y, Inches(0.8), Inches(0.4), ver, font_size=12, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.0), y, Inches(1.5), Inches(0.4), name, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(3.5), y, Inches(8.5), Inches(0.4), desc, font_size=12, color=TEXT_MUTED)

# ════════════════════════════════════════════════════════════════
# 第13页：致谢
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# 装饰圆
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-1.5), Inches(7), Inches(7))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x12, 0x10, 0x25)
circle.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2), "感谢聆听", font_size=48, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5.5), Inches(3.2), Inches(2.3), ACCENT_PURPLE)

add_text_box(slide, Inches(2), Inches(3.6), Inches(9), Inches(0.6), "星尘 Stradust — AI 伴侣应用", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(4.4), Inches(9), Inches(0.5), "Tauri 2 + React 19 + Rust  |  iOS / PC 跨平台", font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.5), "欢迎提问与交流", font_size=20, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = r"F:\stradust\星尘Stradust_演示答辩.pptx"
prs.save(output_path)
print(f"PPT 已保存到: {output_path}")
print(f"共 {len(prs.slides)} 页")
